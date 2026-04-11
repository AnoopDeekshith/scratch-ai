"""
Scratch.ai — Proof of Concept
==============================
Takes a Zoom lecture transcript + lecture slides text and produces
structured study notes in two modes: Detailed and Simplified.

Setup:
    pip install anthropic python-pptx PyPDF2
    export ANTHROPIC_API_KEY="your-api-key-here"

Usage:
    python scratch_ai_poc.py --transcript lecture.vtt --slides lecture.pdf
    python scratch_ai_poc.py --transcript lecture.vtt --slides lecture.pdf --mode simple
    python scratch_ai_poc.py --transcript lecture.txt --slides slides.pptx --mode detailed
"""

import argparse
import os
import sys
import json
from pathlib import Path

# ---------------------------------------------------------------------------
# 1. INPUT PARSERS — handle different file formats
# ---------------------------------------------------------------------------

def parse_vtt_transcript(file_path: str) -> str:
    """Parse a .vtt (WebVTT) transcript file from Zoom into plain text."""
    lines = Path(file_path).read_text(encoding="utf-8").splitlines()
    text_lines = []
    for line in lines:
        line = line.strip()
        # Skip VTT headers, timestamps, and blank lines
        if (
            not line
            or line.startswith("WEBVTT")
            or line.startswith("NOTE")
            or "-->" in line
            or line.isdigit()
        ):
            continue
        text_lines.append(line)
    return " ".join(text_lines)


def parse_txt_transcript(file_path: str) -> str:
    """Read a plain text transcript."""
    return Path(file_path).read_text(encoding="utf-8")


def parse_transcript(file_path: str) -> str:
    """Auto-detect transcript format and parse."""
    ext = Path(file_path).suffix.lower()
    if ext == ".vtt":
        return parse_vtt_transcript(file_path)
    elif ext in (".txt", ".text", ".md"):
        return parse_txt_transcript(file_path)
    else:
        # Try plain text as fallback
        return parse_txt_transcript(file_path)


def parse_pdf_slides(file_path: str) -> str:
    """Extract text from a PDF file (lecture slides)."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        print("ERROR: PyPDF2 not installed. Run: pip install PyPDF2")
        sys.exit(1)

    reader = PdfReader(file_path)
    all_text = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            all_text.append(f"--- Slide/Page {i + 1} ---\n{text}")
    return "\n\n".join(all_text)


def parse_pptx_slides(file_path: str) -> str:
    """Extract text from a PowerPoint file."""
    try:
        from pptx import Presentation
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx")
        sys.exit(1)

    prs = Presentation(file_path)
    all_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            all_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
    return "\n\n".join(all_text)


def parse_slides(file_path: str) -> str:
    """Auto-detect slide format and parse."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return parse_pdf_slides(file_path)
    elif ext in (".pptx", ".ppt"):
        return parse_pptx_slides(file_path)
    elif ext in (".txt", ".md"):
        return Path(file_path).read_text(encoding="utf-8")
    else:
        return Path(file_path).read_text(encoding="utf-8")


# ---------------------------------------------------------------------------
# 2. AI PROCESSING — the core of Scratch.ai
# ---------------------------------------------------------------------------

SYSTEM_PROMPT_DETAILED = """You are Scratch.ai, an AI study assistant that transforms lecture transcripts and slides into comprehensive, well-structured study notes.

Your job is to combine what the teacher SAID (transcript) with what the teacher SHOWED (slides) into notes that are better than either alone.

Output format:

# [Lecture Topic]

## Key Concepts
List the 3-5 most important concepts covered, each with a clear 2-3 sentence explanation.

## Structured Notes
Organize the lecture content by sub-topic. For each sub-topic:
- Use clear headings
- Explain the concept thoroughly
- Include any examples the teacher gave
- Connect it to the slide content
- Highlight definitions, formulas, or key terms in **bold**

## Visual Summary
Describe a simple diagram or flowchart (in text/ASCII) that captures the main relationships between concepts. Use arrows (→) and boxes to show how ideas connect.

## Review Questions
Generate 5 questions a student could use to test their understanding. Mix factual recall and deeper thinking questions.

## Quick Summary
A 3-4 sentence paragraph capturing the entire lecture in a nutshell.

## To-Do / Action Items
List any assignments, readings, deadlines, or tasks the teacher mentioned during the lecture. If none were mentioned, write "No action items mentioned in this lecture."

Be thorough, accurate, and organized. Use the transcript for context and nuance that isn't on the slides. Use the slides for structure and visual content that wasn't spoken aloud."""

SYSTEM_PROMPT_SIMPLE = """You are Scratch.ai, an AI study assistant that transforms lecture transcripts and slides into simple, easy-to-understand study notes written in plain language.

Your job is to take complex lecture material and explain it like you're talking to a smart friend who has zero background in this subject.

Rules:
- No jargon without explanation
- Use analogies and real-world comparisons
- Short sentences
- If a concept is complex, break it into smaller pieces
- Use "think of it like..." and "in other words..." frequently

Output format:

# [Lecture Topic] — The Simple Version

## What Was This Lecture About?
2-3 sentences a non-expert could understand.

## The Big Ideas (Explained Simply)
For each major concept:
- **What it is** — plain English, no jargon
- **Why it matters** — why should you care?
- **Think of it like...** — a real-world analogy

## How It All Connects
A simple explanation of how the concepts relate to each other. Use a basic text diagram if helpful.

## Test Yourself
5 simple questions to check if you actually understood (not just memorized).

## TL;DR
The entire lecture in 2-3 sentences that anyone could understand.

## Don't Forget
Any homework, readings, or deadlines mentioned. If none, write "Nothing mentioned — but review these notes before next class!"

Write like you're explaining to a friend over coffee. Casual but accurate."""


def generate_notes(transcript: str, slides: str, mode: str = "detailed") -> str:
    """Send transcript + slides to Claude API and get structured notes."""
    try:
        import anthropic
    except ImportError:
        print("ERROR: anthropic package not installed. Run: pip install anthropic")
        sys.exit(1)

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("ERROR: ANTHROPIC_API_KEY environment variable not set.")
        print("Get your key at: https://console.anthropic.com/")
        sys.exit(1)

    client = anthropic.Anthropic(api_key=api_key)

    system_prompt = SYSTEM_PROMPT_DETAILED if mode == "detailed" else SYSTEM_PROMPT_SIMPLE

    user_message = f"""Here is the lecture data. Please transform this into structured study notes.

## LECTURE TRANSCRIPT (what the teacher said):
{transcript[:15000]}

## LECTURE SLIDES (what was shown on screen):
{slides[:10000]}

Please generate the study notes now."""

    print(f"\n🧠 Generating {mode} notes with Claude...\n")

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4096,
        system=system_prompt,
        messages=[{"role": "user", "content": user_message}],
    )

    return response.content[0].text


# ---------------------------------------------------------------------------
# 3. OUTPUT — save the results
# ---------------------------------------------------------------------------

def save_notes(notes: str, output_path: str):
    """Save generated notes to a markdown file."""
    Path(output_path).write_text(notes, encoding="utf-8")
    print(f"✅ Notes saved to: {output_path}")


# ---------------------------------------------------------------------------
# 4. MAIN — tie it all together
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Scratch.ai POC — Transform lectures into study notes"
    )
    parser.add_argument(
        "--transcript", "-t",
        required=True,
        help="Path to the lecture transcript (.vtt, .txt)"
    )
    parser.add_argument(
        "--slides", "-s",
        required=True,
        help="Path to the lecture slides (.pdf, .pptx, .txt)"
    )
    parser.add_argument(
        "--mode", "-m",
        choices=["detailed", "simple"],
        default="detailed",
        help="Note style: 'detailed' (default) or 'simple' (layman terms)"
    )
    parser.add_argument(
        "--output", "-o",
        default=None,
        help="Output file path (default: scratch_notes_<mode>.md)"
    )
    parser.add_argument(
        "--both", "-b",
        action="store_true",
        help="Generate both detailed AND simple notes"
    )

    args = parser.parse_args()

    # Validate inputs
    if not Path(args.transcript).exists():
        print(f"ERROR: Transcript file not found: {args.transcript}")
        sys.exit(1)
    if not Path(args.slides).exists():
        print(f"ERROR: Slides file not found: {args.slides}")
        sys.exit(1)

    # Parse inputs
    print("📄 Parsing transcript...")
    transcript = parse_transcript(args.transcript)
    print(f"   → {len(transcript.split())} words extracted")

    print("📊 Parsing slides...")
    slides = parse_slides(args.slides)
    print(f"   → {len(slides.split())} words extracted")

    if not transcript.strip():
        print("WARNING: Transcript appears empty. Check your file.")
    if not slides.strip():
        print("WARNING: Slides appear empty. Check your file.")

    # Generate notes
    if args.both:
        # Generate both modes
        for mode in ["detailed", "simple"]:
            notes = generate_notes(transcript, slides, mode)
            output_path = args.output or f"scratch_notes_{mode}.md"
            if mode == "simple" and args.output:
                # Add suffix for second file
                p = Path(args.output)
                output_path = str(p.parent / f"{p.stem}_simple{p.suffix}")
            save_notes(notes, output_path)
    else:
        notes = generate_notes(transcript, slides, args.mode)
        output_path = args.output or f"scratch_notes_{args.mode}.md"
        save_notes(notes, output_path)

    print("\n🎓 Done! Open the markdown file to review your Scratch.ai notes.")
    print("   Tip: Try both --mode detailed and --mode simple to compare.")


if __name__ == "__main__":
    main()
